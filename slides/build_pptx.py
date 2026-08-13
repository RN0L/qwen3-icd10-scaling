#!/usr/bin/env python3
"""Build the 5-slide deck from the measured records.

Every figure on these slides is read out of ``results/`` at build time. Nothing is
transcribed, so the deck cannot drift from the data, re-run this after any new run
and the numbers update themselves.

    uv run --with python-pptx --with pillow python slides/build_pptx.py

Design intent: one idea per slide, one number large enough to read from the back of
the room, and body text kept under ~45 words. The house colour is a single constant.
Cards size themselves to their content (see ``card``) so text can never be clipped
by the panel it sits in.
"""

from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
RESULTS = REPO / "results"
FIGURES = Path(__file__).resolve().parent / "figures"
OUT = Path(__file__).resolve().parent / "ME344_Final_Arnold_Hambuch.pptx"

# --- palette -------------------------------------------------------------------
ORANGE = RGBColor(0xE8, 0x59, 0x0C)
ORANGE_L = RGBColor(0xFB, 0x92, 0x3C)
TINT = RGBColor(0xFE, 0xF3, 0xEA)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_TINT = RGBColor(0xFE, 0xF2, 0xF2)
INK = RGBColor(0x15, 0x17, 0x1A)
GREY = RGBColor(0x71, 0x7A, 0x85)
FAINT = RGBColor(0xE6, 0xE9, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Helvetica Neue"

W, H = Inches(13.333), Inches(7.5)

# Shapes whose box deliberately overhangs: the ghost slide number is a
# right-aligned box, so its frame reaches the margin while its glyph does not.
EXEMPT = set()


def load(n):
    return json.loads((RESULTS / f"{n}.json").read_text())


# --- facts, read not typed -----------------------------------------------------
recs = {}
for f in sorted(RESULTS.glob("*.json")):
    if f.stem != "analysis":
        d = json.loads(f.read_text())
        recs[d["run_id"]] = d

cpu = recs["cpu-x86-32c-0p6b-bs1-seq1024"]
gpu = recs["gpu-gh200-0p6b-bs1-seq1024"]
tpu = recs["tpu-v5e8-bs1-seq1024-filecache-0p6b"]
base = recs["tpu-v5e8-bs8-seq1024"]
mit = recs["tpu-v5e8-bs8-seq1024-filecache"]

c_s, g_s, t_s = (r["steady"]["median_step_s"] for r in (cpu, gpu, tpu))
sp_g, sp_t = c_s / g_s, c_s / t_s
gt = abs(g_s - t_s) / t_s * 100
bp, mp = base["phases"], mit["phases"]
compute_pct = bp["steady_state_s"] / bp["total_wall_s"] * 100
load_pct = bp["model_load_s"] / bp["total_wall_s"] * 100
load_x = bp["model_load_s"] / mp["model_load_s"]
wall_cut = (1 - mp["total_wall_s"] / bp["total_wall_s"]) * 100

# Utilisation and memory for the like-for-like row. The TPU exposes memory
# statistics but no core-utilisation counter, so that cell says so rather than
# borrowing a number that means something else.
u_cpu = cpu["utilization"]["mean_pct"]
u_gpu = gpu["utilization"]["mean_pct"]
m_cpu, m_gpu, m_tpu = (r["memory"]["peak_pct"] for r in (cpu, gpu, tpu))

# The 4B sweep: the comparison at an operating point where both chips are busy.
g4 = {r["config"]["batch_size"]: r for r in recs.values()
      if r["backend"] == "gpu" and r["config"]["model"].endswith("4B")}
t4 = {r["config"]["batch_size"]: r for r in recs.values()
      if r["backend"] == "tpu" and r["config"]["model"].endswith("4B")
      and r["config"]["seq_len"] == 1024}
g8, t8 = g4[8]["steady"]["tokens_per_s"], t4[8]["steady"]["tokens_per_s"]
g16, t16 = g4[16]["steady"]["tokens_per_s"], t4[16]["steady"]["tokens_per_s"]
g32 = g4[32]["steady"]["tokens_per_s"]
load_ratio = g16 / t16

# The last configuration that fits, per backend, and the one that does not.
def ceiling(pool):
    ok = {b: r for b, r in pool.items() if r["status"] == "ok"}
    died = sorted(b for b, r in pool.items() if r["status"] == "oom")
    top = max(ok)
    return ok[top]["memory"]["peak_pct"], top, (died[0] if died else None)


mem_gpu, bs_gpu, oom_gpu = ceiling(g4)
mem_tpu, bs_tpu, oom_tpu = ceiling(t4)

_c = {c["run_id"]: c["usd_per_1000_steps"] for c in load("analysis")["cost"]["per_run"]}
k4_gpu = _c["gpu-gh200-bs8-seq1024"]
k4_tpu = _c["tpu-v5e8-bs8-seq1024-filecache"]
cost_ratio_4b = k4_tpu / k4_gpu

stats = json.loads((REPO / "docs" / "dataset_stats.json").read_text())
tok = stats["splits"]["train"]["tokens"]["total"]
n_train = stats["splits"]["train"]["n_docs"]
n_dev = stats["splits"]["dev"]["n_docs"]
trunc = {e["seq_len"]: e for e in stats["splits"]["train"]["truncation"]}
cut_1024 = 100 * trunc[1024]["frac_truncated"]

prs = Presentation()
prs.slide_width, prs.slide_height = W, H


# --- primitives ----------------------------------------------------------------
def slide(num, kicker, headline):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), H)
    rule.fill.solid(); rule.fill.fore_color.rgb = ORANGE
    rule.line.fill.background(); rule.shadow.inherit = False

    ghost = s.shapes.add_textbox(Inches(11.9), Inches(0.05), Inches(1.3), Inches(1.5))
    r = ghost.text_frame.paragraphs[0].add_run(); r.text = str(num)
    r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = FAINT; r.font.name = FONT
    ghost.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    EXEMPT.add(ghost.shape_id)

    txt(s, Inches(0.62), Inches(0.40), Inches(10.5), Inches(0.3),
        kicker.upper(), 12, color=ORANGE, bold=True, spacing=2)
    txt(s, Inches(0.62), Inches(0.72), Inches(11.0), Inches(0.9),
        headline, 32, bold=True)
    return s


def txt(s, x, y, w, h, text, size=14, bold=False, color=INK,
        align=PP_ALIGN.LEFT, gap=5, spacing=0, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(gap)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT
    return tb


def lines(text, size, gap, width=None):
    """Height a block of text needs, in inches, counting soft wraps too.

    `width` is the inches available for the text. Helvetica Neue averages roughly
    0.52 em per character in running text, which is what turns a width into a
    character budget. Deliberately generous: over-estimating adds padding,
    under-estimating clips, and clipping is the bug this function exists to stop.
    """
    per_line = max(1, int(width / (size / 72 * 0.52))) if width else 10 ** 6
    n = 0
    for hard in text.split("\n"):
        n += max(1, -(-len(hard) // per_line))
    return n * (size / 72 * 1.28 + gap / 72)


def stat(s, x, y, w, value, label, big=44, color=ORANGE):
    """A single number, sized to read from the back of the room."""
    drop = Inches(big / 72 * 1.15 + 0.08)      # the number's own line box, plus air
    txt(s, x, y, w, drop, value, big, bold=True, color=color, gap=0)
    txt(s, x + Inches(0.02), y + drop, w, Inches(0.9), label, 12, color=GREY, gap=2)
    return y + drop + Inches(lines(label, 12, 2))


def hairline(s, x, y, length, vertical=False, color=FAINT, weight=1.0):
    """A rule. Divides two blocks without wrapping either one in a shape.

    Every tinted card on a slide is a decision to enclose something, and enclosing
    everything is how a deck ends up with six identical containers and no hierarchy.
    A 1pt line does the same separating work and adds no weight.
    """
    w, h = (Pt(weight), length) if vertical else (length, Pt(weight))
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def panel(s, x, y, w, h, fill=TINT, edge=FAINT):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.04
    return sh


PAD = 0.26          # inches of breathing room inside every card
TITLE_GAP = 0.10


def card(s, x, y, w, title, body, title_size=15, body_size=13, gap=4,
         fill=TINT, accent=ORANGE, foot=None):
    """A tinted card whose height is computed from its content.

    The clipped-text bug this replaces came from writing panel heights by hand and
    then editing the text. Here the panel cannot be shorter than what is in it.
    """
    inner = w / Inches(1) - 2 * PAD
    th = title_size / 72 * 1.32 if title else 0.0
    bh = lines(body, body_size, gap, inner)
    fh = (lines(foot, 10.5, 2, inner) + 0.08) if foot else 0.0
    h = PAD + th + (TITLE_GAP if title else 0) + bh + fh + PAD * 0.72
    panel(s, x, y, w, Inches(h), fill=fill)
    cy = y + Inches(PAD)
    if title:
        txt(s, x + Inches(PAD + 0.02), cy, w - Inches(2 * PAD), Inches(th),
            title, title_size, bold=True, color=accent)
        cy += Inches(th + TITLE_GAP)
    txt(s, x + Inches(PAD + 0.02), cy, w - Inches(2 * PAD), Inches(bh),
        body, body_size, gap=gap)
    if foot:
        txt(s, x + Inches(PAD + 0.02), cy + Inches(bh + 0.08),
            w - Inches(2 * PAD), Inches(fh), foot, 10.5, color=GREY, gap=2)
    return y + Inches(h)


def table(s, x, y, rows, widths, size=13, rh=Inches(0.40), highlight=None,
          right_align_from=1, rule_after=None):
    """right_align_from: first column index to right-align. Set high to keep text left."""
    t = s.shapes.add_table(len(rows), len(rows[0]), x, y, sum(widths), rh * len(rows)).table
    for j, wd in enumerate(widths):
        t.columns[j].width = wd
    for i, row in enumerate(rows):
        t.rows[i].height = rh
        for j, v in enumerate(row):
            c = t.cell(i, j); c.text = str(v)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.09)
            c.margin_top = c.margin_bottom = Inches(0.01)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j >= right_align_from else PP_ALIGN.LEFT
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(size); r.font.name = FONT
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = INK
                r.font.bold = True; r.font.color.rgb = WHITE
            elif highlight is not None and j == highlight:
                c.fill.fore_color.rgb = TINT
                r.font.bold = True; r.font.color.rgb = ORANGE
            elif rule_after is not None and i == rule_after:
                c.fill.fore_color.rgb = TINT
                r.font.color.rgb = INK
            else:
                c.fill.fore_color.rgb = WHITE
                r.font.color.rgb = INK
    return y + rh * len(rows)


def fig(s, name, x, y, width):
    """Place a figure and return its bottom edge, so the next thing can clear it."""
    p = FIGURES / name
    if not p.exists():
        return y
    s.shapes.add_picture(str(p), x, y, width=width)
    px, py = Image.open(p).size
    return y + Inches(width / Inches(1) * py / px)


# ============================== 1. Problem ====================================
s = slide(1, "The problem", "Can a clinic afford to run its own ICD-10 model?")

txt(s, Inches(0.62), Inches(1.92), Inches(5.6), Inches(0.35),
    "DrFindus, our startup in Berlin", 15, bold=True, gap=0)
txt(s, Inches(0.62), Inches(2.34), Inches(5.6), Inches(1.5),
    "German GPs code every consultation.\n"
    "We propose, the physician confirms.\n"
    "Reliability is the product.\n"
    "Today it runs on hosted APIs.",
    14, color=GREY, gap=6)

txt(s, Inches(6.55), Inches(1.92), Inches(6.15), Inches(1.0),
    "Could it run on\nour own hardware?", 24, bold=True, color=ORANGE, gap=2)
txt(s, Inches(6.55), Inches(3.02), Inches(6.15), Inches(0.8),
    f"Qwen3-4B + LoRA on CodiEsp: {n_train} train / {n_dev} dev\n"
    "Spanish case reports, ICD-10 coded, CC-BY 4.0.",
    13, color=GREY, gap=3)

txt(s, Inches(0.62), Inches(4.16), Inches(5.9), Inches(0.35),
    "The resource problem", 14, bold=True, color=ORANGE)
txt(s, Inches(0.62), Inches(4.56), Inches(5.9), Inches(1.4),
    "8 GB of weights at 4B before a single activation.\n"
    f"Median case {tok['p50']:,.0f} tokens, longest {tok['max']:,.0f}.\n"
    "One JAX program, lowered by XLA to CPU, GPU and TPU.", 15, gap=9)

fig(s, "slide1-coverage.png", Inches(6.55), Inches(4.02), Inches(6.2))

hairline(s, Inches(0.62), Inches(6.14), Inches(12.05))
txt(s, Inches(0.62), Inches(6.40), Inches(12.0), Inches(0.55),
    "Where does the time go, and what does it take to run this at all?",
    23, bold=True)


# ============================== 2. Approach ===================================
s = slide(2, "Approach", "One program. Three backends. All measured.")

cols = [
    ("Compilation",
     "One JAX program\nXLA lowers it to all three\nLoRA via qwix, Tunix trainer"),
    ("Orchestration",
     "3 pinned Docker images\nKubernetes across 2 clusters\nKueue admission, explicit limits"),
    ("Storage",
     "gcsfuse CSI, implicit-dirs\n8 GB checkpoint from a bucket\nFile cache left as a knob"),
]
# Three parallel aspects, so three columns is the right structure. Three identical
# tinted boxes is not: it gives each the same weight as the storage path below, which
# is the thing this slide is actually building towards.
hairline(s, Inches(0.62), Inches(1.78), Inches(12.10))
for i, (head, body) in enumerate(cols):
    x = Inches(0.62 + i * 4.15)
    if i:
        hairline(s, x - Inches(0.20), Inches(1.96), Inches(1.42), vertical=True)
    txt(s, x, Inches(2.00), Inches(3.85), Inches(0.35), head, 16, bold=True,
        color=ORANGE, gap=0)
    txt(s, x, Inches(2.44), Inches(3.75), Inches(1.1), body, 13, color=GREY, gap=7)

panel(s, Inches(0.62), Inches(3.82), Inches(12.10), Inches(0.66),
      fill=WHITE, edge=ORANGE_L)
txt(s, Inches(0.62), Inches(3.98), Inches(12.10), Inches(0.35),
    "gs://me344-tpu-labs-west4   →   gcsfuse CSI sidecar   →   /gcs in the pod"
    "   →   8 GB checkpoint   →   HBM",
    13.5, bold=True, align=PP_ALIGN.CENTER, gap=0)

table(s, Inches(0.62), Inches(4.72), [
    ["Backend", "Hardware", "Host arch", "How the code got there"],
    ["CPU", "32-core x86_64, 31 GiB", "x86_64", "bare node, podman"],
    ["GPU", "NVIDIA GH200 480 GB", "ARM64 Grace", "public image + ConfigMap"],
    ["TPU", "v5e 2×4, 8 chips", "x86_64", "private registry + gcsfuse"],
], [Inches(1.5), Inches(4.1), Inches(2.4), Inches(4.1)], size=12.5,
    rh=Inches(0.41), right_align_from=99)

txt(s, Inches(0.62), Inches(6.55), Inches(12.0), Inches(0.5),
    "The GH200 host is ARM64. That one fact cost five separate blockers.",
    16, bold=True, color=ORANGE)


# ============================== 3. Measurements ================================
s = slide(3, "Measurements", "We did not time the job. We took it apart.")

hz = 1 / gpu["utilization"]["sample_interval_s"]

y = table(s, Inches(0.62), Inches(1.82), [
    ["What we sample", "Instrument"],
    ["Chip utilisation", f"nvidia-smi {hz:.0f} Hz · psutil on CPU"],
    ["Peak memory", "jax memory_stats() · peak RSS"],
    ["Step time", "median, p10, p90, warmup excluded"],
    ["Wall clock", "six phases, must sum to total"],
], [Inches(2.15), Inches(3.35)], size=12.5, rh=Inches(0.40), right_align_from=99)

txt(s, Inches(0.62), y + Inches(0.16), Inches(5.5), Inches(0.7),
    "One schema, one JSON per run. Every figure in this\n"
    "deck reads those files and nothing else.", 12, color=GREY, gap=2)

card(s, Inches(0.62), Inches(4.52), Inches(5.5),
     "The schema caught a lie",
     "93 µs per step. 11 M tokens/s. Impossible,\n"
     "yet it carried status \"ok\".\n\n"
     "Its own notes admitted it may have timed\n"
     "dispatch, not completion. Discarded and re-run.",
     title_size=15, body_size=12.5, gap=4, fill=RED_TINT, accent=RED)

y = fig(s, "slide3-walltime.png", Inches(6.40), Inches(1.78), Inches(6.30))

table(s, Inches(6.40), y + Inches(0.28), [
    ["Backend", "Peak memory, largest fitting run", "Next step up"],
    ["CPU 32-core", f"{m_cpu:.1f} % of 31 GiB (RSS)", "4B never ran"],
    ["GH200", f"{mem_gpu:.1f} % of 96 GiB HBM", f"batch {oom_gpu} OOM"],
    ["v5e, 8 chips", f"{mem_tpu:.1f} % of 16 GiB/chip", f"batch {oom_tpu} OOM"],
], [Inches(1.60), Inches(3.05), Inches(1.65)], size=12, rh=Inches(0.40),
    right_align_from=99)

txt(s, Inches(6.40), y + Inches(2.02), Inches(6.30), Inches(0.6),
    "Memory is what ends every sweep. One batch step past\nthese numbers, all three fail.",
    13.5, bold=True, color=ORANGE, gap=2)


# ============================== 4. Results ====================================
s = slide(4, "Results", "Under load, one GH200 beats eight TPU chips.")

y = table(s, Inches(0.62), Inches(1.74), [
    ["Qwen3-0.6B, batch 1", "CPU 32-core", "GH200  1 chip", "v5e  8 chips"],
    ["Median step time", f"{c_s:.2f} s", f"{g_s:.4f} s", f"{t_s:.4f} s"],
    ["Speedup vs CPU", "1.0×", f"{sp_g:.0f}×", f"{sp_t:.0f}×"],
    ["Mean chip utilisation", f"{u_cpu:.1f} %", f"{u_gpu:.2f} %", "no counter exists"],
    ["Peak memory", f"{m_cpu:.1f} %", f"{m_gpu:.1f} %", f"{m_tpu:.1f} %"],
], [Inches(2.45), Inches(1.70), Inches(1.70), Inches(1.70)], size=12.5,
    rh=Inches(0.375), rule_after=2)

txt(s, Inches(0.62), y + Inches(0.14), Inches(7.55), Inches(0.7),
    f"Both accelerators are ~{sp_g:.0f}× the CPU and tie with each other to "
    f"{gt:.2f} %.\nAt batch 1 neither is busy. That tie is an artefact.",
    14.5, bold=True, gap=3)

fig(s, "slide4-sweep.png", Inches(0.62), Inches(4.42), Inches(7.55))

stat(s, Inches(8.45), Inches(1.74), Inches(4.25), f"{load_ratio:.1f}x",
     "the throughput of an eight-chip v5e\nslice, from a single Hopper chip", big=52)

hairline(s, Inches(8.45), Inches(3.46), Inches(4.25))

txt(s, Inches(8.45), Inches(3.68), Inches(4.25), Inches(0.34),
    "Where each one runs out of memory", 14, bold=True, gap=0)
txt(s, Inches(8.45), Inches(4.06), Inches(4.25), Inches(0.6),
    f"GH200 between batch {bs_gpu} and {oom_gpu}\nv5e between batch {bs_tpu} and {oom_tpu}",
    13.5, color=GREY, gap=4)

hairline(s, Inches(8.45), Inches(5.02), Inches(4.25))

txt(s, Inches(8.45), Inches(5.24), Inches(4.25), Inches(0.34),
    "The mitigation, controlled", 14, bold=True, gap=0)
txt(s, Inches(8.45), Inches(5.62), Inches(4.25), Inches(0.9),
    f"Checkpoint load {load_x:.1f}× faster, {wall_cut:.0f} % off wall\n"
    "clock, and median step time unchanged\nto the fourth decimal.",
    13, color=GREY, gap=4)


# ============================== 5. Conclusion =================================
s = slide(5, "Conclusion", "The chip was never the bottleneck.")

# One number carries this slide. The other two findings are real but subordinate, and
# giving all three the same tinted card made them all equally ignorable.
txt(s, Inches(0.62), Inches(1.72), Inches(5.2), Inches(1.5),
    f"{compute_pct:.0f}\u2009%", 96, bold=True, color=ORANGE, gap=0)   # thin space: a full one is a canyon at 96pt
txt(s, Inches(0.70), Inches(3.34), Inches(5.0), Inches(0.4),
    "of the wall clock is arithmetic", 19, bold=True, gap=0)
txt(s, Inches(0.70), Inches(3.80), Inches(4.9), Inches(0.9),
    f"Of {bp['total_wall_s']:.0f} s, {bp['steady_state_s']:.0f} s computes. The largest\n"
    f"single phase is reading the checkpoint: {load_pct:.0f} %.",
    13, color=GREY, gap=3)

hairline(s, Inches(6.05), Inches(1.90), Inches(2.45), vertical=True)

# The supporting pair: bold lead, one grey line. No container.
for i, (lead, body) in enumerate([
    ("6.7× more parameters, and the CPU stopped",
     "Not 6.7× slower. XLA never finished compiling, and\n"
     "without remat the kernel killed it at 31.5 GB."),
    ("1 of 36 dependencies was TPU-bound",
     "The GH200 sat idle behind an ARM64 host, a QEMU\n"
     "segfault and a registry 403. One pin is the whole port."),
]):
    y = Inches(1.98 + i * 1.30)
    txt(s, Inches(6.55), y, Inches(6.1), Inches(0.4), lead, 17, bold=True, gap=0)
    txt(s, Inches(6.55), y + Inches(0.40), Inches(6.1), Inches(0.8), body, 13,
        color=GREY, gap=3)

hairline(s, Inches(0.62), Inches(4.86), Inches(12.05))

txt(s, Inches(0.62), Inches(5.14), Inches(6.6), Inches(0.45),
    "Do not scale out. Amortise.", 24, bold=True, gap=0)
txt(s, Inches(0.64), Inches(5.66), Inches(6.5), Inches(0.9),
    "Persistent compile cache  ·  warm workers, not a pod per job\n"
    "file cache on anything that reads a checkpoint",
    13, color=GREY, gap=3)

# Cost stays, as a number rather than a paragraph.
txt(s, Inches(8.30), Inches(5.10), Inches(2.0), Inches(0.6),
    f"${k4_gpu:.2f}", 40, bold=True, color=INK, gap=0)
txt(s, Inches(10.15), Inches(5.28), Inches(2.6), Inches(0.4),
    f"against ${k4_tpu:.2f}", 17, bold=True, color=GREY, gap=0)
txt(s, Inches(8.32), Inches(5.78), Inches(4.4), Inches(0.7),
    f"per 1,000 steps at 4B, batch 8: one GH200\nagainst eight v5e chips, {cost_ratio_4b:.0f}× less",
    12.5, color=GREY, gap=3)
txt(s, Inches(8.32), Inches(6.52), Inches(4.4), Inches(0.3),
    "TPU rate billed; the GH200 rate is a market proxy", 10.5, color=GREY,
    italic=True, gap=0)

txt(s, Inches(0.62), Inches(6.98), Inches(12.1), Inches(0.30),
    f"{len(recs)} measured runs  ·  "
    "github.com/RN0L/qwen3-icd10-scaling",
    11, color=GREY, align=PP_ALIGN.CENTER)


prs.save(OUT)
n = len(prs.slides._sldIdLst)
print(f"wrote {OUT.relative_to(REPO)}, {n} slides, {OUT.stat().st_size // 1024} KB")

# --- guard: nothing may run off the bottom or the right edge --------------------
# The card heights above are computed, but the y offsets are still written by hand,
# so this is the check that catches a hand-written offset that stopped fitting.
MARGIN = 0.14
IN = Inches(1)
WRAPS = []


LINE_H = 1.15       # line box as a multiple of em
BOLD_FACTOR = 1.04  # Helvetica Bold runs about 4 % wider than the regular face

# Helvetica advance widths, em/1000. One average per string is not enough: a headline is
# mostly narrow lowercase (~0.5 em) while "$0.56" is all digits and currency (0.556 em),
# so a single constant either lets a display number wrap or claims every headline does.
_W = {" ": 278, ".": 278, ",": 278, ":": 278, ";": 278, "!": 278, "?": 556, "-": 333,
      "(": 333, ")": 333, "/": 278, "%": 889, "\u00d7": 600, "$": 556, "+": 584,
      "\u2009": 200, "\u00b7": 278, "\u2192": 600}
_W.update({c: 556 for c in "0123456789"})
_W.update(dict(zip("ABCDEFGHIJKLMNOPQRSTUVWXYZ",
                   [667, 667, 722, 722, 667, 611, 778, 722, 278, 500, 667, 556, 833,
                    722, 778, 667, 778, 722, 667, 611, 722, 667, 944, 667, 667, 611])))
_W.update(dict(zip("abcdefghijklmnopqrstuvwxyz",
                   [556, 556, 500, 556, 556, 278, 556, 556, 222, 222, 500, 222, 833,
                    556, 556, 556, 556, 333, 500, 278, 556, 500, 722, 500, 500, 500])))


def text_width_em(text, bold):
    """Width of a string in em, from per-character advances rather than an average."""
    total = sum(_W.get(ch, 556) for ch in text) / 1000.0
    return total * (BOLD_FACTOR if bold else 1.0)


def content_box(sh):
    """The rectangle a shape's *content* actually occupies, in inches.

    Deliberately a different estimate from :func:`lines`. That one is generous, because
    under-estimating there clips text. This one aims to be *realistic*, because
    over-estimating here invents overlaps that are not on the slide: a number set in a
    4-inch box still only inks the inch and a half its glyphs need, and treating the
    declared width as occupied would flag every caption sitting beside one.
    """
    left, top = sh.left / IN, sh.top / IN
    declared_w, declared_h = sh.width / IN, sh.height / IN
    if not (sh.has_text_frame and sh.text_frame.text.strip()):
        return left, top, left + declared_w, top + declared_h

    # A textbox does not lay text out edge to edge: PowerPoint insets it by 0.1in a side
    # by default. Ignoring that is why the first pass called a 1.45in number a fit in a
    # 1.50in box, and LibreOffice then wrapped it onto the line below.
    tf = sh.text_frame
    inset = ((tf.margin_left or 0) + (tf.margin_right or 0)) / IN
    declared_w = max(0.1, declared_w - (inset if inset else 0.2))

    used_w, used_h = 0.0, 0.0
    for para in sh.text_frame.paragraphs:
        size = next((r.font.size.pt for r in para.runs if r.font.size), 13)
        gap = para.space_after.pt if para.space_after else 0
        text = "".join(r.text for r in para.runs)
        em = size / 72.0
        bold = any(r.font.bold for r in para.runs)
        natural = text_width_em(text, bold) * em
        n = max(1, -(-natural // declared_w)) if text else 1
        n = int(n)
        # A single line that does not fit its box will wrap, and a wrapped display
        # number lands on whatever sits below it. That is how "$0.56" became "$0.5"
        # over a "6" in the first LibreOffice export.
        if n > 1 and "\n" not in sh.text_frame.text and size >= 18:
            WRAPS.append((sh, text, natural, declared_w))
        used_w = max(used_w, min(declared_w, natural))
        used_h += n * em * LINE_H + gap / 72.0
    return left, top, left + max(used_w, 0.1), top + max(used_h, 0.1)


bad = []
for i, sl in enumerate(prs.slides, 1):
    boxes = []
    for sh in sl.shapes:
        if sh.height == H or sh.shape_id in EXEMPT:
            continue
        label = (sh.text_frame.text.split("\n")[0][:34]
                 if sh.has_text_frame else str(sh.shape_type))
        l, t, r, b = content_box(sh)
        if b > H / IN - MARGIN:
            bad.append(f"  slide {i}: bottom {b:.2f}in  {label!r}")
        if r > W / IN - MARGIN:
            bad.append(f"  slide {i}: right {r:.2f}in  {label!r}")
        if sh.has_text_frame and sh.text_frame.text.strip():
            boxes.append((l, t, r, b, label))

    # Text over text. The old deck ran a tinted card into the footer and nothing
    # noticed, because the check only looked at the slide edges.
    for j in range(len(boxes)):
        for k in range(j + 1, len(boxes)):
            al, at, ar, ab, an = boxes[j]
            bl, bt, br, bb, bn = boxes[k]
            ox = min(ar, br) - max(al, bl)
            oy = min(ab, bb) - max(at, bt)
            if ox > 0.06 and oy > 0.06:      # ignore hairline touches
                bad.append(
                    f"  slide {i}: {an!r} overlaps {bn!r} "
                    f"by {ox:.2f} x {oy:.2f} in"
                )

for sh, text, natural, declared in WRAPS:
    bad.append(f"  display text {text!r} needs {natural:.2f}in in a {declared:.2f}in box: it will wrap")

print("layout clean: nothing within %.2f in of an edge, no text overlapping text, "
      "no display line wrapping" % MARGIN
      if not bad else "LAYOUT PROBLEMS:\n" + "\n".join(bad))

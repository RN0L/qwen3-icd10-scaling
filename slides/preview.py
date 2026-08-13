#!/usr/bin/env python3
"""Rasterise the deck to PNGs so a layout change can be looked at before it ships.

There is no headless PPTX renderer on this machine — no LibreOffice, and driving Keynote
by AppleScript blocks on a permissions prompt. That left slide design being changed without
anyone seeing the result, which for a deck is the whole of the work. This reads the built
`.pptx` and draws it with Pillow: real positions, real sizes, real text, real colours.

It is an *approximation*, and it is honest about which parts:

* line breaking is greedy on measured text width, close to but not identical to PowerPoint's
* autofit and kerning pairs are not modelled
* pictures are drawn from the source file, so figures are exact
* glyphs the loaded font face lacks come out as boxes here and are fine in the deck —
  U+2192 in the storage path is the one that does this, verified against the real PDF

That is enough to answer the questions a preview has to answer — is anything colliding, is
the hierarchy right, does the slide balance — and not enough to replace the real renderer for
final output.

::

    python3 slides/preview.py                    # -> slides/preview/slide-N.png
    python3 slides/preview.py --scale 2          # denser raster

Requires Pillow and python-pptx, which building the deck already needs.
"""

from __future__ import annotations

import argparse
import os
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
DECK = HERE / "ME344_Final_Arnold_Hambuch.pptx"
OUT = HERE / "preview"

#: Helvetica Neue is what the deck asks for; the .ttc carries the weights in one file.
FONT_FILE = "/System/Library/Fonts/HelveticaNeue.ttc"
FALLBACK = "/System/Library/Fonts/Supplemental/Arial.ttf"
FALLBACK_BOLD = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"

EMU_PER_IN = 914400


def load_font(size_px: int, bold: bool):
    for path, index in ((FONT_FILE, 1 if bold else 0),):
        try:
            return ImageFont.truetype(path, size_px, index=index)
        except Exception:
            pass
    try:
        return ImageFont.truetype(FALLBACK_BOLD if bold else FALLBACK, size_px)
    except Exception:
        return ImageFont.load_default()


def rgb(colour, default=(21, 23, 26)):
    try:
        return tuple(bytes.fromhex(str(colour.rgb)))
    except Exception:
        return default


def wrap(draw, text, font, max_px):
    """Greedy wrap on measured width — the same rule PowerPoint uses, minus kerning."""
    if not text:
        return [""]
    out, line = [], ""
    for word in text.split(" "):
        trial = word if not line else line + " " + word
        if draw.textlength(trial, font=font) <= max_px or not line:
            line = trial
        else:
            out.append(line)
            line = word
    out.append(line)
    return out


def render(slide, index: int, scale: float, width_in: float, height_in: float) -> Path:
    W, H = int(width_in * 96 * scale), int(height_in * 96 * scale)
    img = Image.new("RGB", (W, H), "white")
    draw = ImageDraw.Draw(img)
    px = lambda emu: emu / EMU_PER_IN * 96 * scale

    for shape in slide.shapes:
        left, top = px(shape.left), px(shape.top)
        w, h = px(shape.width), px(shape.height)

        if shape.shape_type == 13 or shape.__class__.__name__ == "Picture":
            try:
                pic = Image.open(io_bytes(shape)).convert("RGB")
                pic = pic.resize((max(1, int(w)), max(1, int(h))))
                img.paste(pic, (int(left), int(top)))
            except Exception:
                draw.rectangle([left, top, left + w, top + h], outline=(200, 205, 210))
            continue

        if shape.has_table:
            tbl = shape.table
            row_tops, acc = [], top
            for row in tbl.rows:
                row_tops.append(acc)
                acc += px(row.height)
            col_lefts, acc = [], left
            for col in tbl.columns:
                col_lefts.append(acc)
                acc += px(col.width)
            for ri, row in enumerate(tbl.rows):
                for ci, cell in enumerate(row.cells):
                    cx, cy = col_lefts[ci], row_tops[ri]
                    cw = px(tbl.columns[ci].width)
                    ch = px(row.height)
                    try:
                        cfill = rgb(cell.fill.fore_color, None) if cell.fill.type == 1 else None
                    except Exception:
                        cfill = None
                    if cfill:
                        draw.rectangle([cx, cy, cx + cw, cy + ch], fill=cfill)
                    para = cell.text_frame.paragraphs[0]
                    runs = para.runs
                    if not runs:
                        continue
                    size_pt = next((r.font.size.pt for r in runs if r.font.size), 12)
                    bold = any(r.font.bold for r in runs)
                    colour = next((rgb(r.font.color) for r in runs if r.font.color
                                   and r.font.color.type is not None), (21, 23, 26))
                    font = load_font(max(1, int(size_pt / 72 * 96 * scale)), bold)
                    text = "".join(r.text for r in runs)
                    tw = draw.textlength(text, font=font)
                    tx = cx + 4 * scale
                    if str(para.alignment) == "RIGHT (3)":
                        tx = cx + cw - tw - 5 * scale
                    elif str(para.alignment) == "CENTER (2)":
                        tx = cx + (cw - tw) / 2
                    draw.text((tx, cy + (ch - size_pt / 72 * 96 * scale) / 2), text,
                              font=font, fill=colour)
            continue

        if shape.has_text_frame and shape.text_frame.text.strip():
            y = top
            for para in shape.text_frame.paragraphs:
                runs = para.runs
                if not runs:
                    y += 6 * scale
                    continue
                size_pt = next((r.font.size.pt for r in runs if r.font.size), 13)
                bold = any(r.font.bold for r in runs)
                colour = next((rgb(r.font.color) for r in runs if r.font.color
                               and r.font.color.type is not None), (21, 23, 26))
                text = "".join(r.text for r in runs)
                font = load_font(max(1, int(size_pt / 72 * 96 * scale)), bold)
                gap = (para.space_after.pt if para.space_after else 0) / 72 * 96 * scale
                for line in wrap(draw, text, font, w):
                    x = left
                    if str(para.alignment) == "CENTER (2)":
                        x = left + (w - draw.textlength(line, font=font)) / 2
                    elif str(para.alignment) == "RIGHT (3)":
                        x = left + w - draw.textlength(line, font=font)
                    draw.text((x, y), line, font=font, fill=colour)
                    y += size_pt / 72 * 96 * scale * 1.18
                y += gap
            continue

        # A plain shape: fill and outline are what carry the panels and rules.
        try:
            fill = rgb(shape.fill.fore_color, None) if shape.fill.type == 1 else None
        except Exception:
            fill = None
        try:
            line = rgb(shape.line.color, None)
        except Exception:
            line = None
        if fill or line:
            draw.rectangle([left, top, left + w, top + h], fill=fill, outline=line)

    OUT.mkdir(exist_ok=True)
    path = OUT / ("slide-%d.png" % index)
    img.save(path)
    return path


def io_bytes(shape):
    import io
    return io.BytesIO(shape.image.blob)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    parser.add_argument("--deck", default=str(DECK))
    parser.add_argument("--scale", type=float, default=1.15)
    args = parser.parse_args()

    if not os.path.exists(args.deck):
        raise SystemExit("%s not found — run `python3 slides/build_pptx.py` first." % args.deck)

    prs = Presentation(args.deck)
    width_in = prs.slide_width / EMU_PER_IN
    height_in = prs.slide_height / EMU_PER_IN
    for i, slide in enumerate(prs.slides, 1):
        print("wrote %s" % render(slide, i, args.scale, width_in, height_in))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
